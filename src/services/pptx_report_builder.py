#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Professional PPTX Report Generator
Based on pptx-generator guidelines with business professional styling
"""

from typing import Dict, List, Optional, Any
from datetime import datetime
from enum import Enum
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None


class PptxReportTheme:
    """Professional color themes for reports"""

    # Corporate Blue - Professional business theme
    CORPORATE = {
        "primary": "1E3A5F",
        "secondary": "4A5568",
        "accent": "D4AF37",
        "light": "E5E7EB",
        "bg": "F9FAFB",
        "text": "1F2937",
        "success": "059669",
        "warning": "D97706",
        "error": "DC2626",
    }

    # Forest Green - Nature/Eco theme
    FOREST = {
        "primary": "2F5233",
        "secondary": "4A6741",
        "accent": "84CC16",
        "light": "ECFCCB",
        "bg": "F7FEE7",
        "text": "1F2937",
        "success": "22C55E",
        "warning": "EAB308",
        "error": "DC2626",
    }

    # Gold Premium - Luxury/Finance
    GOLD = {
        "primary": "B8860B",
        "secondary": "78350F",
        "accent": "FCD34D",
        "light": "FEF3C7",
        "bg": "FFFBEB",
        "text": "1F2937",
        "success": "059669",
        "warning": "D97706",
        "error": "DC2626",
    }

    # Dark Professional
    DARK = {
        "primary": "0F172A",
        "secondary": "334155",
        "accent": "38BDF8",
        "light": "1E293B",
        "bg": "0F172A",
        "text": "F1F5F9",
        "success": "22C55E",
        "warning": "FBBF24",
        "error": "F87171",
    }


class PptxReportBuilder:
    """
    Builder for creating professional PPTX reports
    Following pptx-generator guidelines
    """

    # Slide dimensions: 10" x 5.625" (16:9)
    WIDTH = 10
    HEIGHT = 5.625

    def __init__(self, theme: Dict = None):
        self.theme = theme or PptxReportTheme.CORPORATE
        self.slides = []

    def add_title_slide(
        self, title: str, subtitle: str = "", company: str = "", date: str = None
    ):
        """Add cover/title slide"""
        slide = {
            "type": "cover",
            "title": title,
            "subtitle": subtitle,
            "company": company,
            "date": date or datetime.now().strftime("%Y-%m-%d"),
        }
        self.slides.append(slide)
        return self

    def add_toc_slide(self, sections: List[Dict]):
        """Add table of contents slide"""
        slide = {
            "type": "toc",
            "sections": sections,  # [{"number": "01", "title": "...", "description": "..."}]
        }
        self.slides.append(slide)
        return self

    def add_content_slide(self, title: str, content: Dict):
        """
        Add content slide
        content: {
            "layout": "text" | "two_column" | "chart" | "comparison" | "list",
            "body": ...,
            "chart_data": {...} (optional)
        }
        """
        slide = {"type": "content", "title": title, "content": content}
        self.slides.append(slide)
        return self

    def add_section_divider(self, number: str, title: str):
        """Add section divider slide"""
        slide = {"type": "section", "number": number, "title": title}
        self.slides.append(slide)
        return self

    def add_summary_slide(self, key_points: List[str], next_steps: List[str] = None):
        """Add summary/conclusion slide"""
        slide = {
            "type": "summary",
            "key_points": key_points,
            "next_steps": next_steps or [],
        }
        self.slides.append(slide)
        return self

    def generate_structure(self) -> Dict:
        """Generate report structure for pptx-generator"""
        return {
            "theme": self.theme,
            "slides": self.slides,
            "dimensions": {"width": self.WIDTH, "height": self.HEIGHT},
        }

    def generate_pptx(self, filepath: str) -> str:
        """Generate actual PPTX file"""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx library not available")

        prs = Presentation()
        prs.slide_width = Inches(self.WIDTH)
        prs.slide_height = Inches(self.HEIGHT)

        for slide_data in self.slides:
            slide_type = slide_data.get("type", "content")

            if slide_type == "cover":
                self._add_cover_slide(prs, slide_data)
            elif slide_type == "toc":
                self._add_toc_slide(prs, slide_data)
            elif slide_type == "section":
                self._add_section_slide(prs, slide_data)
            elif slide_type == "summary":
                self._add_summary_slide(prs, slide_data)
            elif slide_type == "content":
                self._add_content_slide(prs, slide_data)

        prs.save(filepath)
        return filepath

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _add_cover_slide(self, prs: Presentation, data: Dict):
        """Add cover slide"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        title = data.get("title", "")
        subtitle = data.get("subtitle", "")
        company = data.get("company", "")
        date = data.get("date", "")

        primary = self._hex_to_rgb(self.theme["primary"])

        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*primary)
                p.alignment = PP_ALIGN.CENTER

    def _add_toc_slide(self, prs: Presentation, data: Dict):
        """Add table of contents slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "المحتويات"

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for section in data.get("sections", []):
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = f"{section.get('number', '')} - {section.get('title', '')}"
            p.font.size = Pt(18)
            p.font.bold = True

    def _add_section_slide(self, prs: Presentation, data: Dict):
        """Add section divider slide"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        number = data.get("number", "")
        title = data.get("title", "")

        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    if not p.text:
                        p.text = f"{number}\n{title}"
                        p.font.size = Pt(44)
                        p.alignment = PP_ALIGN.CENTER

    def _add_summary_slide(self, prs: Presentation, data: Dict):
        """Add summary/conclusion slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = "الخلاصة"

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for point in data.get("key_points", []):
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)

    def _add_content_slide(self, prs: Presentation, data: Dict):
        """Add content slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = data.get("title", "")
        title_shape = slide.shapes.title
        title_shape.text = title

        content = data.get("content", {})
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        layout = content.get("layout", "text")

        if layout == "two_column":
            left = content.get("left", {})
            right = content.get("right", {})

            p = tf.add_paragraph()
            p.text = left.get("title", "")
            p.font.bold = True
            p.font.size = Pt(16)

            for item in left.get("items", []):
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(14)

            p = tf.add_paragraph()
            p.text = "\n" + right.get("title", "")
            p.font.bold = True
            p.font.size = Pt(16)

            for item in right.get("items", []):
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(14)
        else:
            for key, value in content.items():
                if isinstance(value, list):
                    for item in value:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        p.font.size = Pt(14)


class ReportTemplates:
    """Pre-built report templates"""

    @staticmethod
    def sales_report(db_manager, date_range: tuple) -> PptxReportBuilder:
        """Create sales report builder"""
        builder = PptxReportBuilder(PptxReportTheme.CORPORATE)

        # Cover
        builder.add_title_slide(
            "تقرير المبيعات",
            f"الفترة: {date_range[0]} - {date_range[1]}",
            "Logical Version ERP",
        )

        # TOC
        builder.add_toc_slide(
            [
                {
                    "number": "01",
                    "title": "ملخص المبيعات",
                    "description": "إجمالي المبيعات والربحية",
                },
                {
                    "number": "02",
                    "title": "تفاصيل المعاملات",
                    "description": "قائمة الفواتير",
                },
                {
                    "number": "03",
                    "title": "تحليل العملاء",
                    "description": "أفضل العملاء",
                },
                {
                    "number": "04",
                    "title": "المنتجات الأكثر مبيعاً",
                    "description": "أفضل 10 منتجات",
                },
            ]
        )

        # Summary section
        builder.add_section_divider("01", "ملخص المبيعات")

        builder.add_content_slide(
            "ملخص المبيعات",
            {
                "layout": "two_column",
                "left": {
                    "title": "المؤشرات الرئيسية",
                    "items": [
                        "إجمالي المبيعات: 150,000 د.ج",
                        "عدد الفواتير: 45",
                        "متوسط قيمة الفاتورة: 3,333 د.ج",
                        "هامش الربح: 25%",
                    ],
                },
                "right": {
                    "title": "المقارنة مع الفترة السابقة",
                    "items": [
                        "نسبة النمو: +15%",
                        "عدد الفواتير: +8",
                        "متوسط قيمة الفاتورة: +5%",
                        "هامش الربح: +2%",
                    ],
                },
            },
        )

        builder.add_summary_slide(
            ["المبيعات في تزايد مستمر", "هامش الربح يتحسن"],
            ["مراجعة الأسعار", "توسيع خطوط المنتجات"],
        )

        return builder

    @staticmethod
    def inventory_report(db_manager) -> PptxReportBuilder:
        """Create inventory report builder"""
        builder = PptxReportBuilder(PptxReportTheme.FOREST)

        builder.add_title_slide(
            "تقرير المخزون", "حالة المخزون الحالية", "Logical Version ERP"
        )

        builder.add_section_divider("01", "حالة المخزون")

        builder.add_content_slide(
            "ملخص المخزون",
            {
                "layout": "chart",
                "chart_data": {
                    "type": "bar",
                    "title": "حالة المخزون",
                    "data": {"في المخزون": 150, "منخفض": 25, "مرتفع": 30, "نفد": 5},
                },
            },
        )

        builder.add_summary_slide(
            ["المخزون ضمن المستويات الطبيعية", "5 منتجات تحتاج إعادة طلب"],
            ["تحديث حد المخزون الأدنى", "مراجعة دورة الطلب"],
        )

        return builder

    @staticmethod
    def financial_report(db_manager, period: str) -> PptxReportBuilder:
        """Create financial report builder"""
        builder = PptxReportBuilder(PptxReportTheme.GOLD)

        builder.add_title_slide(
            "التقرير المالي", f"الفترة: {period}", "Logical Version ERP"
        )

        builder.add_content_slide(
            "الملخص المالي",
            {
                "layout": "comparison",
                "left": {
                    "title": "الإيرادات",
                    "items": [
                        "المبيعات: 500,000 د.ج",
                        "أخرى: 50,000 د.ج",
                        "الإجمالي: 550,000 د.ج",
                    ],
                },
                "right": {
                    "title": "المصروفات",
                    "items": [
                        "المشتريات: 350,000 د.ج",
                        "التشغيل: 80,000 د.ج",
                        "الإجمالي: 430,000 د.ج",
                    ],
                },
            },
        )

        builder.add_content_slide(
            "صافي الربح",
            {
                "layout": "highlight",
                "value": "120,000 د.ج",
                "label": "صافي الربح",
                "change": "+12%",
            },
        )

        builder.add_summary_slide(
            ["الوضع المالي الصحي", "نمو في الإيرادات"],
            ["خفض المصروفات", "تنويع مصادر الدخل"],
        )

        return builder


class SlideLayouts:
    """Pre-built slide layouts following pptx-generator rules"""

    @staticmethod
    def text_layout(title: str, paragraphs: List[str]) -> Dict:
        """Text-heavy layout"""
        return {"layout": "text", "title": title, "paragraphs": paragraphs}

    @staticmethod
    def two_column_layout(title: str, left: Dict, right: Dict) -> Dict:
        """Two column layout"""
        return {"layout": "two_column", "title": title, "left": left, "right": right}

    @staticmethod
    def chart_layout(title: str, chart_type: str, data: Dict) -> Dict:
        """Chart layout"""
        return {
            "layout": "chart",
            "title": title,
            "chart_type": chart_type,
            "data": data,
        }

    @staticmethod
    def comparison_layout(title: str, option_a: Dict, option_b: Dict) -> Dict:
        """Comparison layout (A vs B)"""
        return {
            "layout": "comparison",
            "title": title,
            "option_a": option_a,
            "option_b": option_b,
        }

    @staticmethod
    def list_layout(title: str, items: List[str], style: str = "bullet") -> Dict:
        """List layout"""
        return {"layout": "list", "title": title, "items": items, "style": style}

    @staticmethod
    def highlight_layout(label: str, value: str, change: str = None) -> Dict:
        """Highlight/stat layout"""
        return {"layout": "highlight", "label": label, "value": value, "change": change}
